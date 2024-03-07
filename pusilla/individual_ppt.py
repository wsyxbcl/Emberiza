from pptx import Presentation
from io import BytesIO
from PIL import Image
import os

def extract_images_with_text_info_from_pptx(pptx_path, output_folder):
    prs = Presentation(pptx_path)

    for i, slide in enumerate(prs.slides):
        # Get the text content and position of each shape on the slide
        text_info = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_info.append({
                    'text': shape.text,
                    'position': (shape.left, shape.top)
                })
        print(text_info)

        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "image"):
                image_stream = BytesIO(shape.image.blob)
                image = Image.open(image_stream)

                # Create a filename using the slide text
                slide_text = "_".join(info['text'] for info in text_info)
                slide_text = ''.join(c for c in slide_text if c.isalnum() or c in [' ', '_', '-']).rstrip()
                
                # Get the position of the image
                image_position = (shape.left, shape.top)

                # Save the image with text and position information in the filename
                filename = f"{output_folder}/{slide_text}_image_{i}_{j}_at_{image_position}.png"
                filename = os.path.normpath(filename)  # Normalize path for different OS
                image.save(filename, format="PNG")
                print(f"Image {i}_{j} saved: {filename}")

# Example usage
pptx_file_path = "/home/wsyxbcl/Downloads/individual_demo.pptx"
output_folder_path = "/home/wsyxbcl/Emberiza/test/images"
extract_images_with_text_info_from_pptx(pptx_file_path, output_folder_path)