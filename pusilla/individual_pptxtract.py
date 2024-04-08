from pptx import Presentation
from io import BytesIO
from PIL import Image
import os
from scipy.spatial import KDTree
from exiftool import ExifToolHelper
import argparse
import pandas as pd

SPECIES = "金钱豹"

TEXT_COORDINATES = {
    (4704845, 136525): f"{SPECIES}名称/姓名标签",
    (7979573, 4433110): "性别",
    (8610805, 5064524): "初次拍摄年/月",
    (10094045, 5708946):"备注",
    (7979574, 5705255): "家族关系",
    (10585259, 4433110): "拍摄地点",
    (11132164, 5071028): "最近捕获年/月",
}

IMG_COORDINATES = {
    (3454708, 1279525): "右侧图",
    (9759950, 1450180): "面部花纹",
    (130654, 1279084): "左侧图",
    (6529770, 1279084): "尾巴",
    (2650811, 4213649): "补充图",
    (130810, 4213860): "前肢花纹",
}

class Individual():
    def __init__(self, species, name, name_label, gender, 
                 first_capture_time, latest_capture_time, 
                 family_relationship, comment, appeared_location):
        self.species = species
        self.name = name
        self.name_label = name_label
        self.gender = gender
        self.first_capture_time = first_capture_time
        self.latest_capture_time = latest_capture_time
        self.family_relationship = family_relationship
        self.comment = comment
        self.appeared_location = appeared_location
        self.age = "adult"

        if "cub" in name_label.lower():
            self.age = "cub"
        elif family_relationship:
            if "'s cub" in family_relationship.lower():
                self.age = "cub"
        elif comment:
            if "'s cub" in comment.lower():
                self.age = "cub"

    def get_description(self):
        first_capture_time = f"初次拍摄于{self.first_capture_time}" if self.first_capture_time else ""
        latest_capture_time = f"最近捕获于{self.latest_capture_time}" if self.latest_capture_time else ""

        # make a list of all the attributes and drop the empty strings
        description = [self.species, self.name, self.name_label, first_capture_time, latest_capture_time, 
                       self.family_relationship, self.comment, self.appeared_location]
        description = [i for i in description if i]
        description_str = (";".join(description)).replace("\n", "")
        return description_str
        
class IndividualImage():
    def __init__(self, individual, body_part, location):
        self.individual = individual
        self.body_part = body_part
        self.location = location

    @property
    def title(self):
        if self.individual.name:
            return f"{self.location}-{self.individual.name}-{self.individual.name_label}"
        else:
            return f"{self.location}-{self.individual.name_label}"

    @property
    def subject(self):
        return self.individual.name_label

    @property
    def keywords(self):
        return [self.individual.species, self.individual.name, self.individual.name_label, self.individual.age, self.body_part]

    @property
    def description(self):
        return self.individual.get_description()

def get_image_name_label(image_position):
    kdtree = KDTree(list(IMG_COORDINATES.keys()))
    _, idx = kdtree.query(image_position)
    return IMG_COORDINATES[list(IMG_COORDINATES.keys())[idx]]

def get_text_info(text_position):
    kdtree = KDTree(list(TEXT_COORDINATES.keys()))
    _, idx = kdtree.query(text_position)
    return TEXT_COORDINATES[list(TEXT_COORDINATES.keys())[idx]]

def patch_text(text):
    '''
    Patch the text to remove human errors and inconsistencies
    '''
    if "错误" in text or "不对" in text:
        return ""
    patched_text = text.replace("\n", "; ")
    patched_text = patched_text.replace("！", "")
    patched_text = patched_text.replace("!", "")
    patched_text = patched_text.replace("‘", "'")
    patched_text = patched_text.replace("的崽子", "'s cub")
    patched_text = patched_text.replace("的娃", "'s cub")
    patched_text = patched_text.replace("club", "cub")
    return patched_text

def extract_images_with_text_info_from_pptx(pptx_path, output_dir, location, info_mode=False):
    prs = Presentation(pptx_path)
    individual_iist = []
    image_list = []
    for i, slide in enumerate(prs.slides):
        # Get the text content and position of each shape on the slide
        text_info = {}
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if shape.text == "":
                    continue
                text_position = (shape.left, shape.top)
                # try:
                #     text_info[TEXT_COORDINATES[text_position]] = shape.text
                # except KeyError:
                #     if shape.text:
                #         print("Unknown TEXT_COORDINATES for {} in slide {}".format(shape.text, i))
                text_info[get_text_info(text_position)] = patch_text(shape.text)
        print(text_info)

        if f"{SPECIES}名称/姓名标签" not in text_info:
            print(f"Slide {i} does not contain {SPECIES}名称/姓名标签")
            continue

        # initialize variables for the individual
        name = gender = first_capture_time = latest_capture_time = family_relationship = comment = appeared_location = None
        try:
            name, name_label = text_info[f"{SPECIES}名称/姓名标签"].split("/")
        except ValueError:
            name_label = text_info[f"{SPECIES}名称/姓名标签"]
            name_label = name_label.strip(" ")
        if "性别" in text_info:
            gender = text_info["性别"]
            if gender in ["M", "公"]:
                gender = "雄"
            if gender in ["F", "母"]:
                gender = "雌"
        if "初次拍摄年/月" in text_info:
            first_capture_time = text_info["初次拍摄年/月"]
        if "最近捕获年/月" in text_info:
            latest_capture_time = text_info["最近捕获年/月"]
        if "家族关系" in text_info:
            family_relationship = text_info["家族关系"]
        if "备注" in text_info:
            comment = text_info["备注"]
        if "拍摄地点" in text_info:
            appeared_location = text_info["拍摄地点"]
        
        individual = Individual(SPECIES, name, name_label, gender, first_capture_time, latest_capture_time, 
                                family_relationship, comment, appeared_location)
        if info_mode:
            individual_iist.append(individual)
        else: 
        # Extract images      
            for j, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    image_stream = BytesIO(shape.image.blob)
                    image = Image.open(image_stream).convert("RGB")

                    # Get the position of the image
                    image_position = (shape.left, shape.top)
                    image_individual = IndividualImage(individual, get_image_name_label(image_position), location)
                    image_list.append(image_individual)
                    # Save the image and write info to metadata
                    # individual_dir = os.path.join(output_dir, individual.name_label)
                    # if not os.path.exists(individual_dir):
                    #     os.makedirs(individual_dir)
                    if not os.path.exists(output_dir):
                        os.makedirs(output_dir)
                    image_path = os.path.join(output_dir, f"{image_individual.location}-{image_individual.subject}-{j}-{image_individual.body_part}.jpg")
                    image.save(image_path, format="JPEG", quality=80)
                    print(f"Saved {image_path}")
                    with ExifToolHelper() as et:
                        et.execute(f"-Subject={image_individual.subject}", 
                                    " ".join([f"-Keywords={keyword}" for keyword in image_individual.keywords]),
                                    f"-Description={image_individual.description}",
                                    f"-Title={image_individual.title}",
                                    "-overwrite_original", 
                                    image_path)

    if info_mode:
        # write individual info to csv
        individual_df = pd.DataFrame([i.__dict__ for i in individual_iist])
        output_csv_path = os.path.join(output_dir, f"{location}-{SPECIES}-individuals.csv")
        individual_df.to_csv(output_csv_path, index=False)
    else:
        # write image info to csv
        individual_df = pd.DataFrame([i.individual.__dict__ for i in image_list])
        image_df = pd.DataFrame([i.__dict__ for i in image_list])
        df = pd.concat([individual_df, image_df], axis=1).drop(columns=["individual"])
        output_csv_path = os.path.join(output_dir, f"{location}-{SPECIES}-images.csv")
        df.to_csv(output_csv_path, index=False)

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract individual info from a pptx file")
    parser.add_argument("pptx_file_path", help="Path to the pptx file")
    parser.add_argument("output_folder_path", help="Path to the output folder")
    parser.add_argument("--info", help="Just output the individual info", action="store_true")
    parser.add_argument('--location', help='Location name', required=True)
    args = parser.parse_args()
    if args.info:
        extract_images_with_text_info_from_pptx(args.pptx_file_path, args.output_folder_path, info_mode=True, location=args.location)
    else:
        extract_images_with_text_info_from_pptx(args.pptx_file_path, args.output_folder_path, location=args.location)
